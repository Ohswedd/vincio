"""Renderers: a :class:`DocumentModel` into a shippable artifact.

Markdown and HTML render dependency-free; DOCX, PDF, and PPTX render behind
opt-in extras (``vincio[gen-docx|gen-pdf|gen-pptx]``). Every renderer consumes
the same IR, so the structural contract is enforced once, before any format is
produced, and the formats never diverge in content.
"""

from __future__ import annotations

import html
import io
from typing import Literal

from pydantic import BaseModel, Field

from ..core.errors import GenerationError
from ..stability import deprecated
from .model import DocBlock, DocumentModel

__all__ = ["RenderFormat", "DocumentArtifact", "render", "MEDIA_TYPES"]

RenderFormat = Literal["markdown", "html", "docx", "pdf", "pptx"]

MEDIA_TYPES: dict[str, str] = {
    "markdown": "text/markdown",
    "html": "text/html",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_TEXT_FORMATS = {"markdown", "html"}

# URL schemes permitted for an image src/target. Anything with an unlisted
# scheme (notably ``javascript:`` / ``vbscript:`` / ``data:text``) is dropped so
# a hostile image path can't smuggle script into rendered HTML/Markdown.
_SAFE_IMAGE_SCHEMES = ("http://", "https://", "data:image/", "file://", "/", "./", "../")


def _safe_image_src(path: str) -> str:
    """A sanitized image src: allowlisted scheme, control chars stripped."""
    cleaned = path.strip().replace("\n", "").replace("\r", "")
    lowered = cleaned.lower()
    if ":" in lowered.split("/", 1)[0]:  # has an explicit scheme
        if not lowered.startswith(_SAFE_IMAGE_SCHEMES):
            return ""
    return cleaned


def _md_image_target(path: str) -> str:
    """A CommonMark-safe image destination — angle-bracketed, breakout-proof."""
    src = _safe_image_src(path)
    if not src:
        return ""
    return "<" + src.replace("<", "%3C").replace(">", "%3E").replace(" ", "%20") + ">"


class DocumentArtifact(BaseModel):
    """A rendered document: its bytes, format, and media type."""

    model_config = {"arbitrary_types_allowed": True}

    format: RenderFormat
    media_type: str
    content: bytes
    title: str = ""
    source_evidence_ids: list[str] = Field(default_factory=list)
    metadata: dict = Field(default_factory=dict)

    @property
    def text(self) -> str:
        """Decoded text for the textual formats (Markdown/HTML)."""
        if self.format not in _TEXT_FORMATS:
            raise GenerationError(f"{self.format!r} is a binary artifact; use .content")
        return self.content.decode("utf-8")

    def save(self, path: str) -> str:
        """Write the artifact to ``path`` and return the path."""
        from pathlib import Path

        Path(path).write_bytes(self.content)
        return path

    def digest(self) -> str:
        """SHA-256 content hash of the rendered bytes."""
        from ..core.media import media_sha256

        return media_sha256(self.content)

    @deprecated(since="7.5", removed_in="8.0", alternative="digest()")
    def sha256(self) -> str:
        """Deprecated name for :meth:`digest`."""
        return self.digest()


# -- table normalization ------------------------------------------------------


def _normalize_table(table) -> tuple[list[str], list[list[str]]]:  # TableData
    """Rectangular ``(columns, rows)`` so every renderer agrees on content.

    Computes the width as ``max(len(columns), widest row)``, synthesizes
    ``colN`` headers when columns are missing, and pads/truncates every row to
    that width — so a ragged or header-less table renders identically (and never
    loses cells) across Markdown/HTML/DOCX/PDF.
    """
    width = max([len(table.columns)] + [len(r) for r in table.rows], default=0)
    if width == 0:
        return [], []
    cols = [str(c) for c in table.columns]
    if len(cols) < width:
        cols += [f"col{i + 1}" for i in range(len(cols), width)]
    rows = [[str(c) for c in row][:width] + [""] * (width - len(row)) for row in table.rows]
    return cols, rows


# -- Markdown -----------------------------------------------------------------


def _md_table(block: DocBlock) -> str:
    table = block.table
    if table is None:
        return ""
    cols, rows = _normalize_table(table)
    if not cols:
        return ""
    lines: list[str] = []
    if table.title:
        lines.append(f"**{table.title}**")
        lines.append("")
    lines.append("| " + " | ".join(c.replace("|", r"\|") for c in cols) + " |")
    lines.append("| " + " | ".join("---" for _ in cols) + " |")
    for row in rows:
        lines.append("| " + " | ".join(str(c).replace("|", r"\|") for c in row) + " |")
    for note in table.footnotes:
        lines.append("")
        lines.append(f"_{note}_")
    return "\n".join(lines)


def render_markdown(model: DocumentModel) -> str:
    parts: list[str] = []
    if model.title:
        parts.append(f"# {model.title}")
    if model.subtitle:
        parts.append(f"*{model.subtitle}*")
    for block in model.blocks:
        if block.kind == "heading":
            level = min(6, max(1, block.level + 1))  # title is h1; sections start at h2
            parts.append(f"{'#' * level} {block.text}")
        elif block.kind == "paragraph":
            parts.append(block.text)
        elif block.kind == "quote":
            parts.append("\n".join(f"> {line}" for line in block.text.splitlines() or [""]))
        elif block.kind == "code":
            parts.append(f"```{block.language}\n{block.text}\n```")
        elif block.kind == "list":
            marker = (lambda i: f"{i + 1}.") if block.ordered else (lambda i: "-")
            parts.append("\n".join(f"{marker(i)} {item}" for i, item in enumerate(block.items)))
        elif block.kind == "table":
            parts.append(_md_table(block))
        elif block.kind == "image" and block.image_path:
            target = _md_image_target(block.image_path)
            alt = (block.image_alt or block.text).replace("[", "(").replace("]", ")")
            if target:
                line = f"![{alt}]({target})"
                parts.append(line + (f"\n\n*{block.text}*" if block.text else ""))
        elif block.kind == "page_break":
            parts.append("---")
    if model.footnotes:
        parts.append("## Notes")
        parts.append("\n".join(f"{i + 1}. {note}" for i, note in enumerate(model.footnotes)))
    if model.bibliography:
        parts.append("## Sources")
        parts.append("\n".join(f"- {entry}" for entry in model.bibliography))
    return "\n\n".join(p for p in parts if p).strip() + "\n"


# -- HTML ---------------------------------------------------------------------


def _html_table(block: DocBlock) -> str:
    table = block.table
    if table is None:
        return ""
    cols, rows = _normalize_table(table)
    if not cols:
        return ""
    out = ["<table>"]
    if table.title:
        out.append(f"<caption>{html.escape(table.title)}</caption>")
    out.append("<thead><tr>")
    out.extend(f"<th>{html.escape(c)}</th>" for c in cols)
    out.append("</tr></thead>")
    out.append("<tbody>")
    for row in rows:
        out.append("<tr>" + "".join(f"<td>{html.escape(c)}</td>" for c in row) + "</tr>")
    out.append("</tbody></table>")
    return "".join(out)


def render_html(model: DocumentModel) -> str:
    body: list[str] = []
    if model.title:
        body.append(f"<h1>{html.escape(model.title)}</h1>")
    if model.subtitle:
        body.append(f'<p class="subtitle"><em>{html.escape(model.subtitle)}</em></p>')
    for block in model.blocks:
        if block.kind == "heading":
            level = min(6, max(2, block.level + 1))
            anchor = f' id="{html.escape(block.anchor)}"' if block.anchor else ""
            body.append(f"<h{level}{anchor}>{html.escape(block.text)}</h{level}>")
        elif block.kind == "paragraph":
            body.append(f"<p>{html.escape(block.text)}</p>")
        elif block.kind == "quote":
            body.append(f"<blockquote>{html.escape(block.text)}</blockquote>")
        elif block.kind == "code":
            body.append(f"<pre><code>{html.escape(block.text)}</code></pre>")
        elif block.kind == "list":
            tag = "ol" if block.ordered else "ul"
            items = "".join(f"<li>{html.escape(i)}</li>" for i in block.items)
            body.append(f"<{tag}>{items}</{tag}>")
        elif block.kind == "table":
            body.append(_html_table(block))
        elif block.kind == "image" and block.image_path:
            src = _safe_image_src(block.image_path)
            if src:
                cap = f"<figcaption>{html.escape(block.text)}</figcaption>" if block.text else ""
                alt = html.escape(block.image_alt or block.text)
                body.append(f'<figure><img src="{html.escape(src)}" alt="{alt}">{cap}</figure>')
        elif block.kind == "page_break":
            body.append('<hr class="page-break">')
    if model.footnotes:
        body.append("<section class='notes'><h2>Notes</h2><ol>")
        body.extend(f"<li>{html.escape(n)}</li>" for n in model.footnotes)
        body.append("</ol></section>")
    if model.bibliography:
        body.append("<section class='sources'><h2>Sources</h2><ul>")
        body.extend(f"<li>{html.escape(e)}</li>" for e in model.bibliography)
        body.append("</ul></section>")
    title = html.escape(model.title or "Document")
    return (
        "<!DOCTYPE html>\n<html lang=\"en\"><head><meta charset=\"utf-8\">"
        f"<title>{title}</title></head>\n<body>\n" + "\n".join(body) + "\n</body></html>\n"
    )


# -- DOCX (python-docx) -------------------------------------------------------


def render_docx(model: DocumentModel) -> bytes:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - optional dep
        raise GenerationError('DOCX generation requires python-docx: pip install "vincio[gen-docx]"') from exc

    document = docx.Document()
    if model.title:
        document.add_heading(model.title, level=0)
    if model.subtitle:
        document.add_paragraph(model.subtitle, style="Intense Quote")
    for block in model.blocks:
        if block.kind == "heading":
            document.add_heading(block.text, level=min(9, max(1, block.level)))
        elif block.kind == "paragraph":
            document.add_paragraph(block.text)
        elif block.kind == "quote":
            document.add_paragraph(block.text, style="Quote")
        elif block.kind == "code":
            document.add_paragraph(block.text)
        elif block.kind == "list":
            style = "List Number" if block.ordered else "List Bullet"
            for item in block.items:
                document.add_paragraph(item, style=style)
        elif block.kind == "table" and block.table is not None:
            cols, rows = _normalize_table(block.table)
            if cols:
                docx_table = document.add_table(rows=1, cols=len(cols))
                docx_table.style = "Light Grid Accent 1"
                for i, col in enumerate(cols):
                    docx_table.rows[0].cells[i].text = col
                for row in rows:
                    cells = docx_table.add_row().cells
                    for i, value in enumerate(row):
                        cells[i].text = value
        elif block.kind == "image" and block.image_path:
            try:
                document.add_picture(block.image_path)
            except Exception:  # noqa: BLE001 - missing/unsupported image file
                document.add_paragraph(f"[image: {block.image_path}]")
            if block.text:
                document.add_paragraph(block.text, style="Caption")
        elif block.kind == "page_break":
            document.add_page_break()
    if model.footnotes:
        document.add_heading("Notes", level=1)
        for i, note in enumerate(model.footnotes, start=1):
            document.add_paragraph(f"{i}. {note}")
    if model.bibliography:
        document.add_heading("Sources", level=1)
        for entry in model.bibliography:
            document.add_paragraph(entry, style="List Bullet")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


# -- PDF (reportlab) ----------------------------------------------------------


def render_pdf(model: DocumentModel) -> bytes:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import (
            ListFlowable,
            ListItem,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            TableStyle,
        )
        from reportlab.platypus import Table as PdfTable
    except ImportError as exc:  # pragma: no cover - optional dep
        raise GenerationError('PDF generation requires reportlab: pip install "vincio[gen-pdf]"') from exc

    styles = getSampleStyleSheet()
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=LETTER, title=model.title or "Document")
    flow: list = []

    def esc(text: str) -> str:
        return html.escape(text)

    if model.title:
        flow.append(Paragraph(esc(model.title), styles["Title"]))
    if model.subtitle:
        flow.append(Paragraph(esc(model.subtitle), styles["Italic"]))
        flow.append(Spacer(1, 12))
    for block in model.blocks:
        if block.kind == "heading":
            style = styles[f"Heading{min(4, max(1, block.level))}"]
            flow.append(Paragraph(esc(block.text), style))
        elif block.kind in ("paragraph", "quote"):
            flow.append(Paragraph(esc(block.text), styles["BodyText"]))
        elif block.kind == "code":
            flow.append(Paragraph(esc(block.text), styles["Code"]))
        elif block.kind == "list":
            items = [ListItem(Paragraph(esc(i), styles["BodyText"])) for i in block.items]
            flow.append(ListFlowable(items, bulletType="1" if block.ordered else "bullet"))
        elif block.kind == "table" and block.table is not None:
            cols, rows = _normalize_table(block.table)
            if cols:
                data = [cols] + rows
                pdf_table = PdfTable(data, repeatRows=1)
                pdf_table.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2f5496")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                            ("FONTSIZE", (0, 0), (-1, -1), 8),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ]
                    )
                )
                flow.append(pdf_table)
        elif block.kind == "image" and block.image_path:
            from reportlab.platypus import Image as PdfImage

            try:
                flow.append(PdfImage(block.image_path, width=400, height=300, kind="proportional"))
            except Exception:  # noqa: BLE001
                flow.append(Paragraph(esc(f"[image: {block.image_path}]"), styles["BodyText"]))
            if block.text:
                flow.append(Paragraph(esc(block.text), styles["Italic"]))
        elif block.kind == "page_break":
            from reportlab.platypus import PageBreak

            flow.append(PageBreak())
        flow.append(Spacer(1, 6))
    if model.footnotes:
        flow.append(Paragraph("Notes", styles["Heading1"]))
        for i, note in enumerate(model.footnotes, start=1):
            flow.append(Paragraph(esc(f"{i}. {note}"), styles["BodyText"]))
    if model.bibliography:
        flow.append(Paragraph("Sources", styles["Heading1"]))
        for entry in model.bibliography:
            flow.append(Paragraph(esc(entry), styles["BodyText"]))
    doc.build(flow)
    return buffer.getvalue()


# -- PPTX (python-pptx) -------------------------------------------------------


def render_pptx(model: DocumentModel) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - optional dep
        raise GenerationError('PPTX generation requires python-pptx: pip install "vincio[gen-pptx]"') from exc

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = model.title or "Presentation"
    if model.subtitle and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = model.subtitle

    bullet_layout = prs.slide_layouts[1]
    for head, body in model.sections():
        if not (head.text or body):
            continue
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = head.text or model.title or "Section"
        frame = slide.placeholders[1].text_frame
        frame.clear()
        first = True
        for block in body:
            lines: list[str] = []
            if block.kind in ("paragraph", "quote", "code") and block.text:
                lines.append(block.text)
            elif block.kind == "list":
                lines.extend(block.items)
            elif block.kind == "table" and block.table is not None:
                lines.append(block.table.to_text())
            elif block.kind == "image" and block.image_path:
                try:
                    slide.shapes.add_picture(block.image_path, Inches(5), Inches(2), height=Inches(3))
                except Exception:  # noqa: BLE001
                    lines.append(f"[image: {block.image_path}]")
            for line in lines:
                paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                paragraph.text = line
                paragraph.font.size = Pt(18)
                first = False
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


_RENDERERS = {
    "markdown": lambda m: render_markdown(m).encode("utf-8"),
    "html": lambda m: render_html(m).encode("utf-8"),
    "docx": render_docx,
    "pdf": render_pdf,
    "pptx": render_pptx,
}


def render(model: DocumentModel, fmt: RenderFormat) -> DocumentArtifact:
    """Render ``model`` into a :class:`DocumentArtifact` in format ``fmt``."""
    if fmt not in _RENDERERS:
        raise GenerationError(f"unknown render format {fmt!r}; known: {sorted(_RENDERERS)}")
    content = _RENDERERS[fmt](model)
    return DocumentArtifact(
        format=fmt,
        media_type=MEDIA_TYPES[fmt],
        content=content,
        title=model.title,
        source_evidence_ids=list(model.source_evidence_ids),
        metadata=dict(model.metadata),
    )
